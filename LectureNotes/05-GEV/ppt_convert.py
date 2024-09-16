from bs4 import BeautifulSoup
from pptx import Presentation

# 打开HTML文件
with open('05slides.html', 'r', encoding='utf-8') as file:
    html_content = file.read()

# 使用BeautifulSoup解析HTML
soup = BeautifulSoup(html_content, 'html.parser')

# 查找包含幻灯片内容的textarea
textarea = soup.find('textarea', {'id': 'source'})

# 如果找到textarea，按'---'分隔幻灯片内容
if textarea:
    slides_content = textarea.text.split('---')

    # 创建PPT对象
    prs = Presentation()

    # 遍历每一张幻灯片内容
    for slide_text in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 使用标题+文本布局
        title, content = slide_text.strip().split('\n', 1)  # 分离标题和内容

        # 添加标题
        slide.shapes.title.text = title.strip()

        # 添加内容到文本框
        text_box = slide.shapes.placeholders[1]
        text_box.text = content.strip()

    # 保存PPT文件
    prs.save('output.pptx')

    print("PPT文件已生成：output.pptx")
else:
    print("没有找到幻灯片内容")
